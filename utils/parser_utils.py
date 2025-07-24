from docx import Document
from pptx import Presentation
import pandas as pd
import fitz

def parse_files(files):
    chunks = []
    for file in files:
        if file.name.endswith(".pdf"):
            chunks += parse_pdf(file)
        elif file.name.endswith(".docx"):
            chunks += parse_docx(file)
        elif file.name.endswith(".pptx"):
            chunks += parse_pptx(file)
        elif file.name.endswith(".csv"):
            chunks += parse_csv(file)
        elif file.name.endswith(".txt"):
            chunks += file.read().decode("utf-8").split("\n")
    return chunks

def parse_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return [page.get_text() for page in doc]

def parse_docx(file):
    doc = Document(file)
    return [para.text for para in doc.paragraphs if para.text.strip()]

def parse_pptx(file):
    prs = Presentation(file)
    slides = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        slides.append(text)
    return slides

def parse_csv(file):
    df = pd.read_csv(file)
    return df.astype(str).apply(lambda row: " ".join(row), axis=1).tolist()
